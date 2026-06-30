"""PowerPoint (.pptx) status report generator (python-pptx)."""
import io
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PRIMARY = RGBColor(0x25, 0x63, 0xEB)
DARK = RGBColor(0x0F, 0x17, 0x2A)


def _title_slide(prs, project, customer, week):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(8.8), Inches(2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = f"Weekly Status Report\n{project}"
    run.font.size = Pt(34); run.font.bold = True; run.font.color.rgb = DARK
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = f"{customer}  |  Week ending {week}  |  SD Advisory"
    r2.font.size = Pt(14); r2.font.color.rgb = PRIMARY


def _bullet_slide(prs, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    head = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.9))
    hr = head.text_frame.paragraphs[0].add_run()
    hr.text = title; hr.font.size = Pt(26); hr.font.bold = True; hr.font.color.rgb = PRIMARY

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    tf = body.text_frame; tf.word_wrap = True
    items = items or ["None."]
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = f"•  {it}"
        r.font.size = Pt(16); r.font.color.rgb = DARK


def build_status_pptx(report: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)
    _title_slide(prs, report.get("project", ""), report.get("customer", ""),
                 report.get("week", date.today()))
    _bullet_slide(prs, "Executive Summary", [report.get("exec_summary", "")])
    _bullet_slide(prs, "Completed This Period", report.get("completed", []))
    _bullet_slide(prs, "Planned Next Period", report.get("upcoming", []))
    _bullet_slide(prs, "Risks & Issues", (report.get("risks", []) + report.get("issues", [])))
    _bullet_slide(prs, "Budget & Schedule",
                  [report.get("budget_status", ""), report.get("schedule_status", "")])
    _bullet_slide(prs, "Decisions Required", report.get("decisions", []))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
